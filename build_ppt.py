from __future__ import annotations

import argparse
import json
import logging
from pathlib import Path
from typing import Any

from jsonschema import validate
from pptx import Presentation
from pptx.util import Inches

from parser.html_free_layout import parse_html_to_render
from renderer.free_layout import render_absolute_elements
from utils.theme import Theme


logging.basicConfig(level=logging.INFO, format="[%(levelname)s] %(message)s")


def _load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _validate(doc: dict[str, Any], schema_path: Path) -> None:
    schema = _load_json(schema_path)
    validate(instance=doc, schema=schema)


def _ensure_free_layout_slide(pres: Presentation, render_obj: dict[str, Any]) -> None:
    blank = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank)
    render_absolute_elements(slide, render_obj.get("elements", []))


def _presentation(theme: Theme) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(theme.slide_width_inch)
    prs.slide_height = Inches(theme.slide_height_inch)
    return prs


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", type=Path, help="slides JSON spec path")
    parser.add_argument("--html", type=Path, help="single HTML path for free_layout parse")
    parser.add_argument("--output", type=Path, required=True, help="output pptx path")
    parser.add_argument("--theme", type=Path, default=Path("config/theme.json"))
    args = parser.parse_args()

    theme = Theme.load(args.theme)
    prs = _presentation(theme)

    if args.html:
        render_obj = parse_html_to_render(args.html, theme)
        _validate(render_obj, Path("config/render.schema.json"))
        for d in render_obj.get("diagnostics", []):
            logging.warning("%s", d)
        _ensure_free_layout_slide(prs, render_obj)
    elif args.input:
        slides_doc = _load_json(args.input)
        _validate(slides_doc, Path("config/slides.schema.json"))
        for slide in slides_doc.get("slides", []):
            if slide.get("page_type") == "free_layout":
                render_obj = slide.get("render", {})
                _validate(render_obj, Path("config/render.schema.json"))
                for d in render_obj.get("diagnostics", []):
                    logging.warning("slide %s: %s", slide.get("page_id"), d)
                _ensure_free_layout_slide(prs, render_obj)
            else:
                # deterministic placeholder for template mode until dedicated templates are implemented
                _ensure_free_layout_slide(
                    prs,
                    {
                        "mode": "absolute_elements",
                        "elements": [
                            {
                                "id": f"title_{slide.get('page_id')}",
                                "type": "text",
                                "left": 0.8,
                                "top": 0.8,
                                "width": 11.6,
                                "height": 1.0,
                                "text": slide.get("content", {}).get("title", slide.get("page_id", "")),
                                "style": {
                                    "font_family": "Microsoft YaHei",
                                    "font_size_pt": 24,
                                    "bold": True,
                                    "color": [28, 45, 74],
                                    "align": "left",
                                    "valign": "middle",
                                },
                            }
                        ],
                        "diagnostics": ["Template mode placeholder used"],
                    },
                )
    else:
        raise ValueError("Either --input or --html is required")

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    logging.info("Saved PPTX to %s", args.output)


if __name__ == "__main__":
    main()
