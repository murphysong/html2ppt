from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any

from jsonschema import validate
from pptx import Presentation
from pptx.util import Inches

from parser.html_free_layout import parse_html_to_render
from renderer.free_layout import render_absolute_elements
from utils.theme import Theme
import json


def _load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _numeric_sort_key(path: Path) -> tuple[int, str]:
    stem = path.stem
    return (int(stem), stem) if stem.isdigit() else (10**9, stem)


def _new_presentation(theme: Theme) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(theme.slide_width_inch)
    prs.slide_height = Inches(theme.slide_height_inch)
    return prs


def _add_render_slide(prs: Presentation, render_obj: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_absolute_elements(slide, render_obj.get("elements", []))


def main() -> None:
    parser = argparse.ArgumentParser(description="Batch convert examples HTML into individual PPTX and one merged PPTX")
    parser.add_argument("--examples-dir", type=Path, default=Path("examples"))
    parser.add_argument("--output-dir", type=Path, default=Path("output/examples_ppt"))
    parser.add_argument("--merged-output", type=Path, default=Path("output/examples_merged.pptx"))
    parser.add_argument("--theme", type=Path, default=Path("config/theme.json"))
    args = parser.parse_args()

    theme = Theme.load(args.theme)
    render_schema = _load_json(Path("config/render.schema.json"))

    html_files = sorted(args.examples_dir.glob("*.html"), key=_numeric_sort_key)
    if not html_files:
        raise ValueError(f"No html files found in {args.examples_dir}")

    args.output_dir.mkdir(parents=True, exist_ok=True)
    merged_prs = _new_presentation(theme)

    for html in html_files:
        render_obj = parse_html_to_render(html, theme)
        validate(instance=render_obj, schema=render_schema)

        single_prs = _new_presentation(theme)
        _add_render_slide(single_prs, render_obj)
        single_output = args.output_dir / f"{html.stem}.pptx"
        single_prs.save(single_output)

        _add_render_slide(merged_prs, render_obj)
        print(f"[OK] {html.name} -> {single_output}")

    args.merged_output.parent.mkdir(parents=True, exist_ok=True)
    merged_prs.save(args.merged_output)
    print(f"[OK] merged -> {args.merged_output}")
    print(f"[DONE] total slides: {len(html_files)}")


if __name__ == "__main__":
    main()
