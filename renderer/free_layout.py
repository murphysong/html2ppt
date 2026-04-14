from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


SHAPE_MAP = {
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "rounded_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
    "circle": MSO_AUTO_SHAPE_TYPE.OVAL,
    "panel": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "rule_bar": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
}

ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
VALIGN_MAP = {"top": MSO_VERTICAL_ANCHOR.TOP, "middle": MSO_VERTICAL_ANCHOR.MIDDLE, "bottom": MSO_VERTICAL_ANCHOR.BOTTOM}


def _rgb(arr: list[int] | None, default: tuple[int, int, int]) -> RGBColor:
    if not arr:
        arr = [*default]
    return RGBColor(int(arr[0]), int(arr[1]), int(arr[2]))


def render_absolute_elements(slide: Any, elements: list[dict[str, Any]]) -> None:
    for el in sorted(elements, key=lambda x: x.get("z_index", 0)):
        typ = el["type"]
        style = el.get("style", {})
        if typ == "shape":
            shp = slide.shapes.add_shape(
                SHAPE_MAP.get(el.get("shape", "rect"), MSO_AUTO_SHAPE_TYPE.RECTANGLE),
                Inches(el["left"]), Inches(el["top"]), Inches(el["width"]), Inches(el["height"]),
            )
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(style.get("fill_color"), (255, 255, 255))
            shp.line.color.rgb = _rgb(style.get("line_color"), (255, 255, 255))
            shp.line.width = Pt(style.get("line_width_pt", 0.0))
        elif typ == "text":
            tb = slide.shapes.add_textbox(Inches(el["left"]), Inches(el["top"]), Inches(el["width"]), Inches(el["height"]))
            tf = tb.text_frame
            tf.clear()
            tf.vertical_anchor = VALIGN_MAP.get(style.get("valign", "top"), MSO_VERTICAL_ANCHOR.TOP)
            p = tf.paragraphs[0]
            p.alignment = ALIGN_MAP.get(style.get("align", "left"), PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = el.get("text", "")
            run.font.name = style.get("font_family", "Microsoft YaHei")
            run.font.size = Pt(style.get("font_size_pt", 14))
            run.font.bold = bool(style.get("bold", False))
            run.font.italic = bool(style.get("italic", False))
            run.font.color.rgb = _rgb(style.get("color"), (51, 51, 51))
        elif typ == "line":
            ln = slide.shapes.add_connector(
                1,
                Inches(el["x1"]), Inches(el["y1"]), Inches(el["x2"]), Inches(el["y2"]),
            )
            ln.line.color.rgb = _rgb(style.get("line_color"), (200, 200, 200))
            ln.line.width = Pt(style.get("line_width_pt", 1.0))
